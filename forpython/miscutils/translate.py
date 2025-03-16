#!/usr/bin/env python
"""===============================================================================

        FILE: /Users/nailbiter/for/forpython/miscutils/translate.py

       USAGE: (not intended to be directly executed)

 DESCRIPTION: 

     OPTIONS: ---
REQUIREMENTS: ---
        BUGS: ---
       NOTES: ---
      AUTHOR: Alex Leontiev (alozz1991@gmail.com)
ORGANIZATION: 
     VERSION: ---
     CREATED: 2025-03-06T19:23:30.722590
    REVISION: ---

TODO:
  1(D). add tenacity
  2(D). custom translation function
  3(D). sqlite3 cache
==============================================================================="""

import click
from dotenv import load_dotenv
import os
from tenacity import retry, stop_after_attempt, wait_fixed
from os import path
import logging
import typing
import functools
import sys
from pptx import Presentation
from pptx.text.text import TextFrame
import tqdm
import sys
from openpyxl import load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from alex_leontiev_toolbox_python.utils.disk_cache import FsCache
import json
from _handy import load_function_from_file

moption = functools.partial(click.option, show_envvar=True)
KNOWN_FILE_FORMATS = {"pptx": None, "plaintext": None, "xlsx": None}
KNOWN_LANGUAGES = {"es": None, "ja": None, "en": None}


def replace_text_in_xlsx(xlsx_file: str, replace_func: typing.Callable) -> None:
    """
    Iterates through all text elements in an .xlsx file, replaces them with the
    output of a provided Python function, and outputs the modified .xlsx to stdout.

    Args:
        xlsx_file (str or file-like object): Path to the .xlsx file or a file-like object.
        replace_func (callable): A function that takes the original text as input
                                   and returns the replacement text.
    """

    wb = load_workbook(xlsx_file)
    pbar = tqdm.tqdm()
    for sheet in tqdm.tqdm(wb):
        for row in sheet.iter_rows():
            for cell in row:
                pbar.update(1)
                if cell.value is not None and isinstance(cell.value, str):
                    original_text = cell.value
                    replacement_text = replace_func(original_text)
                    cell.value = replacement_text
    pbar.close()

    # Save the modified workbook to stdout
    wb.save(sys.stdout.buffer)


def replace_text_in_pptx(pptx_file, replace_func):
    """
    Iterates through all text elements in a .pptx file, replaces them with the
    output of a provided Python function, and outputs the modified .pptx to stdout.

    Args:
        pptx_file (str or file-like object): Path to the .pptx file or a file-like object.
        replace_func (callable): A function that takes the original text as input
                                   and returns the replacement text.
    """

    prs = Presentation(pptx_file)

    for slide in tqdm.tqdm(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        replacement_text = replace_func(original_text)
                        run.text = replacement_text

    # Save the modified presentation to stdout
    prs.save(sys.stdout.buffer)


@click.command()
@moption("--input-file", "-i", type=click.Path(allow_dash=True), default="-")
@moption(
    "--input-format",
    "-F",
    type=click.Choice(["auto", *KNOWN_FILE_FORMATS]),
    default="auto",
)
@moption(
    "--input-language",
    "-I",
    type=click.Choice(list(KNOWN_LANGUAGES)),
    required=True,
    default="ja",
)
@moption(
    "--output-language",
    "-O",
    type=click.Choice(list(KNOWN_LANGUAGES)),
    required=True,
    default="en",
)
@moption(
    "-M", "--method", type=click.Choice(["googletrans", "official"]), default="official"
)
@moption("-S", "--salt", type=str)
@moption("--custom-translation-callback", "-C", type=(click.Path(), str))
@moption(
    "--translate-source-language/--no-translate-source-language", "-T/ ", default=False
)
@moption("--kwargs-json", "-K", type=str, default="{}")
def translate(
    input_file,
    translate_source_language,
    input_format,
    input_language,
    output_language,
    method,
    custom_translation_callback,
    salt,
    kwargs_json,
):
    logging.warning(
        dict(input_language=input_language, output_language=output_language)
    )

    translate_kwargs = dict(
        json.loads(kwargs_json),
        translate_source_language=translate_source_language,
        method=method,
        input_language=input_language,
        output_language=output_language,
        salt=salt,
    )
    if custom_translation_callback is None:
        translation_callback = translate_text
    else:
        translation_callback = load_function_from_file(*custom_translation_callback)

    if input_format == "auto" or input_format == "plaintext":
        with click.open_file(input_file) as f:
            text = f.read().strip()
        logging.warning(f"i: {text}")
        text = translation_callback(text, **translate_kwargs)

        logging.warning(f"o: {text}")
        click.echo(text)
    elif input_format == "pptx":
        replace_text_in_pptx(
            input_file,
            replace_func=functools.partial(translation_callback, **translate_kwargs),
        )
    elif input_format == "xlsx":
        replace_text_in_xlsx(
            input_file,
            replace_func=functools.partial(translation_callback, **translate_kwargs),
        )
    else:
        raise NotImplementedError(dict(input_format=input_format))


@functools.cache
@FsCache(path.join(path.dirname(__file__), ".translate.cache.d"), is_loud=True)
@retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
def translate_text(
    text: typing.Optional[str],
    method: str,
    input_language: str,
    output_language: str,
    salt: typing.Optional[str] = None,
    translate_source_language: bool = False,
    source_language_detection_confidence: float = 0.1,
    is_accept_none_confidence: bool = True,
) -> typing.Optional[str]:
    if text is None:
        return None
    elif text.strip() == "":
        return text

    if method == "googletrans":
        #! pip install 'googletrans==4.0.0rc1'
        #! pip install googletrans
        from googletrans import Translator

        translator = Translator()
        try:
            should_translate = True
            if not translate_source_language:
                result = translator.detect(text)
                logging.warning(result)
                if result.lang == output_language and (
                    (is_accept_none_confidence and (result.confidence is None))
                    or (result.confidence >= source_language_detection_confidence)
                ):
                    should_translate = False
            logging.warning(f"should_translate: {should_translate}")
            if should_translate:
                text = translator.translate(
                    text, src=input_language, dest=output_language
                ).text
        except Exception as e:
            logging.error(e)
            raise (
                e,
                dict(
                    text=text,
                    method=method,
                    input_language=input_language,
                    output_language=output_language,
                    salt=salt,
                ),
            )
    elif method == "official":
        #! pip install google-cloud-translate
        #! pip install google-cloud-translate==2.0.1
        # from google.cloud import translate as google_translate
        from google.cloud import translate_v2 as google_translate

        client = google_translate.Client()
        result = client.translate(
            text, source_language=input_language, target_language=output_language
        )["translatedText"]
    else:
        raise NotImplementedError(dict(method=method))

    return text


if __name__ == "__main__":
    fn = ".env"
    if path.isfile(fn):
        logging.warning(f"loading `{fn}`")
        load_dotenv(dotenv_path=fn)
    translate(show_default=True, auto_envvar_prefix="TRANSLATE")
